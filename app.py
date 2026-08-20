from fastapi import FastAPI, Request
from fastapi.responses import HTMLResponse, JSONResponse, FileResponse
import google.generativeai as genai
import docx
from pptx import Presentation
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
import os

app = FastAPI()

# HTML Frontend الشامل (القائمة الجانبية، تسجيل الدخول، الدفع والاشتراكات، والذكاء الاصطناعي)
html_content = """
<!DOCTYPE html>
<html lang="ar" dir="rtl">
<head>
    <meta charset="UTF-8">
    <title>NOX AI Studio | منظومة سليمان ماهر المتكاملة</title>
    <style>
        :root {
            --bg-sidebar: #171717;
            --bg-main: #212121;
            --bg-input: #2f2f2f;
            --border-color: #383838;
            --text-main: #ececec;
            --text-muted: #9b9b9b;
            --accent-color: #1a73e8;
            --accent-hover: #1557b0;
        }
        body {
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            background-color: var(--bg-main);
            color: var(--text-main);
            margin: 0;
            display: flex;
            height: 100vh;
            overflow: hidden;
        }
        /* القائمة الجانبية */
        .sidebar {
            width: 300px;
            background-color: var(--bg-sidebar);
            border-left: 1px solid var(--border-color);
            display: flex;
            flex-direction: column;
            padding: 20px;
            box-sizing: border-box;
            justify-content: space-between;
        }
        .sidebar h2 { font-size: 18px; color: #fff; margin-bottom: 20px; text-align: center; }
        .settings-group { margin-bottom: 15px; }
        .settings-group label { font-size: 12px; color: var(--text-muted); display: block; margin-bottom: 5px; }
        .settings-group input {
            width: 100%;
            background: var(--bg-input);
            border: 1px solid var(--border-color);
            color: #fff;
            padding: 8px;
            border-radius: 6px;
            font-size: 13px;
            box-sizing: border-box;
            outline: none;
        }
        .user-panel {
            background: var(--bg-input);
            padding: 12px;
            border-radius: 8px;
            border: 1px solid var(--border-color);
            margin-bottom: 15px;
        }
        .action-btn {
            background-color: var(--accent-color);
            color: white;
            border: none;
            width: 100%;
            padding: 10px;
            border-radius: 6px;
            cursor: pointer;
            font-weight: 600;
            margin-top: 8px;
            font-size: 13px;
        }
        .action-btn:hover { background-color: var(--accent-hover); }
        .sub-btn { background-color: #059669; }
        .sub-btn:hover { background-color: #047857; }

        /* المحتوى الرئيسي */
        .main-content {
            flex-grow: 1;
            display: flex;
            flex-direction: column;
            position: relative;
            height: 100vh;
        }
        .chat-messages {
            flex-grow: 1;
            overflow-y: auto;
            padding: 20px;
            display: flex;
            flex-direction: column;
            gap: 20px;
            max-width: 900px;
            width: 100%;
            margin: 0 auto;
            box-sizing: border-box;
        }
        .message {
            background: #262626;
            padding: 15px;
            border-radius: 12px;
            border: 1px solid var(--border-color);
            line-height: 1.6;
        }
        .input-area {
            padding: 20px;
            background: var(--bg-main);
            display: flex;
            justify-content: center;
        }
        .input-box-wrapper {
            background-color: var(--bg-input);
            border: 1px solid var(--border-color);
            border-radius: 16px;
            padding: 12px 15px;
            width: 100%;
            max-width: 900px;
            display: flex;
            flex-direction: column;
            gap: 10px;
        }
        textarea {
            width: 100%;
            background: transparent;
            border: none;
            color: var(--text-main);
            font-size: 15px;
            outline: none;
            resize: none;
            height: 45px;
            font-family: inherit;
        }
        .send-btn {
            background-color: var(--accent-color);
            color: white;
            border: none;
            padding: 8px 16px;
            border-radius: 8px;
            cursor: pointer;
            font-weight: 600;
            align-self: flex-end;
        }
        .export-buttons {
            display: flex;
            gap: 8px;
            margin-top: 10px;
            flex-wrap: wrap;
        }
        .export-btn {
            background: #334155;
            color: #fff;
            border: none;
            padding: 6px 12px;
            border-radius: 6px;
            font-size: 12px;
            cursor: pointer;
        }
        .export-btn:hover { background: #475569; }

        /* نافذة الاشتراك والدفع المنبثقة (Modal) */
        .modal {
            display: none;
            position: fixed;
            z-index: 1000;
            left: 0; top: 0; width: 100%; height: 100%;
            background-color: rgba(0,0,0,0.7);
            justify-content: center;
            align-items: center;
        }
        .modal-content {
            background-color: var(--bg-sidebar);
            padding: 25px;
            border-radius: 12px;
            border: 1px solid var(--border-color);
            width: 400px;
            text-align: center;
        }
        .pricing-card {
            background: var(--bg-input);
            border: 1px solid var(--border-color);
            padding: 15px;
            border-radius: 8px;
            margin: 10px 0;
            cursor: pointer;
        }
        .pricing-card:hover { border-color: var(--accent-color); }
    </style>
</head>
<body>

    <!-- القائمة الجانبية -->
    <div class="sidebar">
        <div>
            <h2>NOX AI Studio</h2>

            <div class="user-panel">
                <div id="userInfo" style="font-size: 13px; margin-bottom: 8px;">مرحباً: <b>زائر</b></div>
                <button class="action-btn" id="authBtn" onclick="toggleAuth()">تسجيل الدخول</button>
            </div>

            <div class="settings-group">
                <label>مفتاح Gemini API Key</label>
                <input type="password" id="apiKey" placeholder="AIzaSy...">
            </div>

            <div class="user-panel" style="text-align: center;">
                <div style="font-size: 12px; color: var(--text-muted);">حالة الاشتراك</div>
                <div id="subStatus" style="font-size: 14px; font-weight: bold; color: #f59e0b; margin: 5px 0;">باقة مجانية</div>
                <button class="action-btn sub-btn" onclick="openSubModal()">ترقية الاشتراك والدفع 💳</button>
            </div>
        </div>

        <div style="font-size: 12px; color: var(--text-muted); text-align: center; border-top: 1px solid var(--border-color); padding-top: 10px;">
            <b>سليمان ماهر</b> © 2026
        </div>
    </div>

    <!-- المحتوى الرئيسي -->
    <div class="main-content">
        <div class="chat-messages" id="chatMessages">
            <div class="message">
                <strong>أهلاً بك يا سليمان في منظومتك المتكاملة.</strong>
                <p style="color: var(--text-muted); margin: 5px 0 0 0;">القائمة الجانبية جاهزة لإدارة حسابك ومفتاح الـ API والاشتراكات، وابدأ بكتابة طلبك بالأسفل.</p>
            </div>
        </div>

        <div class="input-area">
            <div class="input-box-wrapper">
                <textarea id="promptInput" placeholder="اكتب طلبك هنا (مثال: جدول احترافي اكسل للربع الاول أو تقرير شامل)..." onkeydown="handleKey(event)"></textarea>
                <button class="send-btn" onclick="sendMessage()">إرسال ➔</button>
            </div>
        </div>
    </div>

    <!-- نافذة الاشتراكات والدفع -->
    <div class="modal" id="subModal">
        <div class="modal-content">
            <h3>اختر خطة الاشتراك</h3>
            <p style="font-size: 12px; color: var(--text-muted);">استمتع بمميزات غير محدودة وتصدير احترافي</p>

            <div class="pricing-card" onclick="processPayment('باقة PRO الشهرية - 49 ريال')">
                <h4>باقة PRO</h4>
                <p style="font-size: 13px; color: var(--text-muted);">تصدير غير محدود + نماذج متقدمة</p>
                <b>49 ر.س / شهرياً</b>
            </div>

            <div class="pricing-card" onclick="processPayment('الباقة الذهبية للأعمال - 149 ريال')">
                <h4>باقة الأعمال (Enterprise)</h4>
                <p style="font-size: 13px; color: var(--text-muted);">دعم كامل للشركات والربط البرمجي</p>
                <b>149 ر.س / سنوياً</b>
            </div>

            <button class="action-btn" style="background: #ef4444; margin-top: 15px;" onclick="closeSubModal()">إغلاق</button>
        </div>
    </div>

    <script>
        let isLoggedIn = false;

        function toggleAuth() {
            isLoggedIn = !isLoggedIn;
            document.getElementById('userInfo').innerHTML = isLoggedIn ? 'مرحباً: <b>سليمان ماهر</b>' : 'مرحباً: <b>زائر</b>';
            document.getElementById('authBtn').innerText = isLoggedIn ? 'تسجيل الخروج' : 'تسجيل الدخول';
            document.getElementById('authBtn').style.background = isLoggedIn ? '#374151' : '#1a73e8';
        }

        function openSubModal() { document.getElementById('subModal').style.display = 'flex'; }
        function closeSubModal() { document.getElementById('subModal').style.display = 'none'; }

        function processPayment(planName) {
            alert('تم اختيار: ' .concat(planName, '\\nجاري توجيهك لبوابة الدفع الآمنة...'));
            document.getElementById('subStatus').innerText = 'مشترك (' + planName.split(' ')[0] + ')';
            document.getElementById('subStatus').style.color = '#10b981';
            closeSubModal();
        }

        function handleKey(e) {
            if (e.key === 'Enter' && !e.shiftKey) {
                e.preventDefault();
                sendMessage();
            }
        }

        async function sendMessage() {
            const text = document.getElementById('promptInput').value.trim();
            const apiKey = document.getElementById('apiKey').value;
            const chat = document.getElementById('chatMessages');
            if (!text) return;
            if (!apiKey) { alert('الرجاء إدخال مفتاح API في القائمة الجانبية أولاً'); return; }

            chat.innerHTML += `<div class="message"><b>أنت:</b> ${text}</div>`;
            document.getElementById('promptInput').value = '';

            const res = await fetch('/chat', {
                method: 'POST',
                headers: { 'Content-Type': 'application/json' },
                body: JSON.stringify({ prompt: text, apiKey })
            });
            const data = await res.json();

            chat.innerHTML += `
                <div class="message">
                    <b>AI:</b>
                    <div style="white-space: pre-wrap; margin-top: 8px;">${data.reply}</div>
                    <div class="export-buttons">
                        <button class="export-btn" onclick="exportFile('${encodeURIComponent(text)}', 'excel')">📥 Excel (Q1)</button>
                        <button class="export-btn" onclick="exportFile('${encodeURIComponent(text)}', 'word')">📥 Word</button>
                        <button class="export-btn" onclick="exportFile('${encodeURIComponent(text)}', 'powerpoint')">📥 PowerPoint</button>
                        <button class="export-btn" onclick="exportFile('${encodeURIComponent(text)}', 'pdf')">📥 Text/PDF</button>
                    </div>
                </div>
            `;
            chat.scrollTop = chat.scrollHeight;
        }

        async function exportFile(topic, type) {
            const apiKey = document.getElementById('apiKey').value;
            alert('جاري تجهيز وتنزيل الملف...');
            const res = await fetch('/generate', {
                method: 'POST',
                headers: { 'Content-Type': 'application/json' },
                body: JSON.stringify({ topic: decodeURIComponent(topic), apiKey, serviceType: type })
            });
            if (res.ok) {
                const blob = await res.blob();
                const url = window.URL.createObjectURL(blob);
                const a = document.createElement('a');
                a.href = url;
                const ext = type === 'excel' ? 'xlsx' : type === 'word' ? 'docx' : type === 'powerpoint' ? 'pptx' : 'txt';
                a.download = `Gemini_File.${ext}`;
                document.body.appendChild(a);
                a.click();
                a.remove();
            } else {
                alert('فشل التصدير.');
            }
        }
    </script>
</body>
</html>
"""


def get_best_model(api_key):
    genai.configure(api_key=api_key)
    for m in genai.list_models():
        if 'generateContent' in m.supported_generation_methods and 'flash' in m.name:
            return m.name
    return 'gemini-1.5-flash'


@app.get("/", response_class=HTMLResponse)
async def read_index():
    return html_content


@app.post("/chat")
async def chat_endpoint(request: Request):
    data = await request.json()
    try:
        model = genai.GenerativeModel(get_best_model(data['apiKey']))
        response = model.generate_content(data['prompt'])
        return {"reply": response.text}
    except Exception as e:
        return JSONResponse(content={"detail": str(e)}, status_code=400)


@app.post("/generate")
async def generate_endpoint(request: Request):
    data = await request.json()
    topic = data.get("topic")
    api_key = data.get("apiKey")
    service_type = data.get("serviceType")

    try:
        model = genai.GenerativeModel(get_best_model(api_key))
        ai_content = model.generate_content(f"اكتب تفاصيل وتقريراً احترافيًا حول: {topic}").text
    except:
        ai_content = f"محتوى خاص بـ: {topic}"

    if service_type == "excel":
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "ميزانية الربع الأول Q1"
        ws.sheet_view.rightToLeft = True

        header_fill = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
        header_font = Font(name="Segoe UI", size=11, bold=True, color="FFFFFF")
        regular_font = Font(name="Segoe UI", size=10)
        bold_font = Font(name="Segoe UI", size=10, bold=True)
        thin_border = Border(
            left=Side(style='thin', color='D9D9D9'), right=Side(style='thin', color='D9D9D9'),
            top=Side(style='thin', color='D9D9D9'), bottom=Side(style='thin', color='D9D9D9')
        )

        headers = ["البند / القسم", "يناير", "فبراير", "مارس", "الإجمالي الربعي"]
        ws.append(headers)

        for col_num in range(1, len(headers) + 1):
            cell = ws.cell(row=1, column=col_num)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center")

        q1_rows = [
            ["الإيرادات والمبيعات", 50000, 55000, 60000, "=SUM(B2:D2)"],
            ["التكاليف التشغيلية", 20000, 22000, 21000, "=SUM(B3:D3)"],
            ["الرواتب والأجور", 15000, 15000, 15000, "=SUM(B4:D4)"],
            ["المصروفات النثرية", 5000, 4000, 4500, "=SUM(B5:D5)"],
            ["صافي الربح التشغيلي", "=B2-SUM(B3:B5)", "=C2-SUM(C3:C5)", "=D2-SUM(D3:D5)", "=E2-SUM(E3:E5)"]
        ]

        for row_idx, row_data in enumerate(q1_rows, start=2):
            ws.append(row_data)
            for col_idx in range(1, len(row_data) + 1):
                cell = ws.cell(row=row_idx, column=col_idx)
                cell.font = bold_font if row_idx == 6 else regular_font
                cell.border = thin_border
                cell.alignment = Alignment(horizontal="center", vertical="center")

        for col in ws.columns:
            max_len = max(len(str(cell.value or '')) for cell in col)
            col_letter = get_column_letter(col[0].column)
            ws.column_dimensions[col_letter].width = max(max_len + 5, 18)

        path = "Q1_Professional_Sheet.xlsx"
        wb.save(path)
        return FileResponse(path, filename="Q1_Professional_Sheet.xlsx")

    elif service_type == "word":
        doc = docx.Document()
        doc.add_heading(f'تقرير: {topic}', 0)
        for line in ai_content.split('\n'):
            if line.strip(): doc.add_paragraph(line)
        path = "Gemini_Report.docx"
        doc.save(path)
        return FileResponse(path, filename="Gemini_Report.docx")

    elif service_type == "powerpoint":
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = topic
        slide.placeholders[1].text = "منظومة سليمان ماهر الذكية"

        lines = [line.strip() for line in ai_content.split('\n') if line.strip()]
        chunk_size = 5
        for i in range(0, len(lines), chunk_size):
            chunk = lines[i:i + chunk_size]
            if not chunk: continue
            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            slide2.shapes.title.text = f"تفاصيل: {topic}"
            slide2.placeholders[1].text = "\n".join(chunk)

        path = "Gemini_Presentation.pptx"
        prs.save(path)
        return FileResponse(path, filename="Gemini_Presentation.pptx")

    else:
        path = "Gemini_Document.txt"
        with open(path, "w", encoding="utf-8") as f:
            f.write(ai_content)
        return FileResponse(path, filename="Gemini_Document.txt")